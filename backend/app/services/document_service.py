from typing import Optional, List
import os
from pathlib import Path
import PyPDF2
from pptx import Presentation
from docx import Document as DocxDocument

from app.models.document import Document, ProcessingStatus
from app.config import settings

class DocumentService:
    @staticmethod
    def extract_text_from_file(file_path: str, document_type: str) -> Optional[str]:
        """Extract text content from uploaded file"""
        try:
            if document_type == 'pdf':
                return DocumentService._extract_from_pdf(file_path)
            elif document_type in ['ppt', 'pptx']:
                return DocumentService._extract_from_pptx(file_path)
            elif document_type in ['doc', 'docx']:
                return DocumentService._extract_from_docx(file_path)
            elif document_type == 'txt':
                return DocumentService._extract_from_txt(file_path)
            else:
                return None
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return None
    
    @staticmethod
    def _extract_from_pdf(file_path: str) -> str:
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text.strip()
    
    @staticmethod
    def _extract_from_pptx(file_path: str) -> str:
        text = ""
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    
    @staticmethod
    def _extract_from_docx(file_path: str) -> str:
        doc = DocxDocument(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    
    @staticmethod
    def _extract_from_txt(file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    
    @staticmethod
    def process_document(document: Document) -> bool:
        """Process document and extract content"""
        try:
            # Extract text content
            raw_content = DocumentService.extract_text_from_file(
                document.file_path, 
                document.document_type.value
            )
            
            if raw_content:
                document.raw_content = raw_content
                document.processed_content = DocumentService._clean_text(raw_content)
                document.processing_status = ProcessingStatus.COMPLETED
                
                # TODO: Add more processing like:
                # - Language detection
                # - Summary generation using LLM
                # - Knowledge point extraction
                
                return True
            else:
                document.processing_status = ProcessingStatus.FAILED
                document.processing_error = "Failed to extract text content"
                return False
                
        except Exception as e:
            document.processing_status = ProcessingStatus.FAILED
            document.processing_error = str(e)
            return False
    
    @staticmethod
    def _clean_text(text: str) -> str:
        """Clean and normalize extracted text"""
        # Remove excessive whitespace
        lines = [line.strip() for line in text.split('\n')]
        lines = [line for line in lines if line]  # Remove empty lines
        return '\n'.join(lines)